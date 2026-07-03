"""
build_presentation.py
=====================
Generates chassis_workbench_presentation.pptx
Run: python3 presentation/build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0d, 0x11, 0x17)   # near-black background
C_SURFACE  = RGBColor(0x16, 0x1b, 0x22)   # card surface
C_ACCENT   = RGBColor(0x58, 0xa6, 0xff)   # blue accent
C_ACCENT2  = RGBColor(0x3f, 0xb9, 0x50)   # green
C_CYAN     = RGBColor(0x39, 0xd3, 0xf0)   # cyan
C_WARN     = RGBColor(0xf7, 0x8a, 0x0e)   # amber
C_WHITE    = RGBColor(0xff, 0xff, 0xff)
C_MUTED    = RGBColor(0x8b, 0x94, 0x9e)
C_PURPLE   = RGBColor(0xbc, 0x8c, 0xff)

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(14), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=Pt(13), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def accent_bar(slide, color=C_ACCENT, w=Inches(0.06)):
    add_rect(slide, Inches(0.4), Inches(1.1), w, Inches(5.8), fill=color)


def slide_title(slide, title, subtitle=None, accent=C_ACCENT):
    # Top accent line
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=accent)
    # Bottom accent line
    add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill=accent)

    add_text(slide, title,
             Inches(0.6), Inches(0.18), Inches(11), Inches(0.7),
             size=Pt(32), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(0.82), Inches(11), Inches(0.4),
                 size=Pt(15), color=C_MUTED)


def bullet_box(slide, items, l, t, w, h,
               title=None, title_color=C_ACCENT,
               bullet="▸", item_size=Pt(13), gap=Pt(6),
               bg=None, border=C_SURFACE):
    if bg:
        add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(0.5))
    txb = slide.shapes.add_textbox(l + Inches(0.15), t + Inches(0.12),
                                   w - Inches(0.3), h - Inches(0.2))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = title_color

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = gap
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = item_size
        run.font.color.rgb = C_WHITE
    return txb


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ═════════════════════════════════════════════════════════════════════════════

def slide_hero(prs):
    s = blank_slide(prs)
    fill_bg(s)

    # Full-width gradient-like rectangle (simulated with overlapping rects)
    add_rect(s, Inches(0), Inches(0), W, H, fill=RGBColor(0x0a, 0x0e, 0x14))

    # Accent stripe left
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=C_ACCENT)

    # Big title
    txb = slide.shapes.add_textbox if False else s.shapes.add_textbox(
        Inches(0.55), Inches(1.5), Inches(8.5), Inches(1.8))
    txb.word_wrap = False
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Motorcycle Chassis"
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "Dynamics Workbench"
    r2.font.size = Pt(46)
    r2.font.bold = True
    r2.font.color.rgb = C_ACCENT

    # Subtitle
    add_text(s, "R&D Engineering Platform  ·  Physics-First  ·  Foale / Cossalter Methods",
             Inches(0.6), Inches(3.55), Inches(9), Inches(0.5),
             size=Pt(16), color=C_MUTED)

    # Tag pills
    for i, (label, color) in enumerate([
        ("13 Physics Modules", C_ACCENT),
        ("8 Bike Families",    C_ACCENT2),
        ("Real-time Backend",  C_CYAN),
        ("Validated 86%",      C_WARN),
    ]):
        x = Inches(0.6 + i * 2.8)
        add_rect(s, x, Inches(4.3), Inches(2.55), Inches(0.42),
                 fill=RGBColor(0x1a, 0x22, 0x2e), line=color, line_w=Pt(1.2))
        add_text(s, label, x + Inches(0.1), Inches(4.32),
                 Inches(2.4), Inches(0.38), size=Pt(13), bold=True, color=color)

    # Footer
    add_text(s, "2026  ·  261631710+dikshant1103-beep@users.noreply.github.com",
             Inches(0.6), Inches(7.0), Inches(9), Inches(0.35),
             size=Pt(11), color=C_MUTED)
    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Background
# ═════════════════════════════════════════════════════════════════════════════

def slide_background(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Background", "Why this tool exists")

    # Left column — Problem
    bullet_box(s,
        [
            "Existing tools (MSC ADAMS, VI-Grade, dSPACE) cost $50k–$500k/year",
            "No accessible real-time chassis design tool built for motorcycles",
            "Spreadsheet methods are fragmented — one formula per file",
            "Race engineers rely on tribal knowledge, not live physics feedback",
            "Academic tools lack the full Foale + Cossalter coupled pipeline",
        ],
        Inches(0.4), Inches(1.3), Inches(5.8), Inches(5.7),
        title="The Problem",
        title_color=C_WARN,
        bg=RGBColor(0x18, 0x1a, 0x1f),
        border=C_WARN,
    )

    # Right column — Solution
    bullet_box(s,
        [
            "Open-source Electron desktop app — runs fully offline",
            "Python FastAPI backend: all 13 physics modules unified",
            "TypeScript frontend: instant visual feedback while backend syncs",
            "Based strictly on peer-reviewed Foale & Cossalter formulas",
            "Validated against published manufacturer specs (Yamaha R1, MT-09)",
        ],
        Inches(6.7), Inches(1.3), Inches(6.2), Inches(5.7),
        title="Our Approach",
        title_color=C_ACCENT2,
        bg=RGBColor(0x18, 0x1a, 0x1f),
        border=C_ACCENT2,
    )

    add_rect(s, Inches(6.35), Inches(1.3), Inches(0.04), Inches(5.7), fill=C_SURFACE)
    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AIM
# ═════════════════════════════════════════════════════════════════════════════

def slide_aim(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Aim", "What this software is designed to do")

    goals = [
        ("Primary Goal",
         "Provide motorcycle R&D engineers a single, physics-accurate desktop tool "
         "to design, analyze, and validate chassis geometry and dynamics in real time.",
         C_ACCENT),
        ("Design Target",
         "Replace disconnected spreadsheets and expensive simulation licenses with "
         "an integrated platform that gives instant results backed by real physics.",
         C_ACCENT2),
        ("Validation Target",
         "All outputs validated against Foale Ch.2–11 / Cossalter Ch.1–8 formulas "
         "and published manufacturer data. Current score: 38/44 checks pass (86%).",
         C_CYAN),
        ("Long-term Target",
         "Extend to full multi-body dynamics (MBD), contact mechanics, Pacejka tire "
         "model, and stability eigenvalue analysis for professional race engineering.",
         C_PURPLE),
    ]

    for i, (heading, body, color) in enumerate(goals):
        row = i // 2
        col = i % 2
        x = Inches(0.4 + col * 6.45)
        y = Inches(1.3 + row * 2.85)
        add_rect(s, x, y, Inches(6.1), Inches(2.6),
                 fill=RGBColor(0x16, 0x1c, 0x26), line=color, line_w=Pt(1.5))
        add_rect(s, x, y, Inches(6.1), Inches(0.06), fill=color)
        add_text(s, heading, x + Inches(0.2), y + Inches(0.14),
                 Inches(5.7), Inches(0.4), size=Pt(14), bold=True, color=color)
        add_text(s, body, x + Inches(0.2), y + Inches(0.58),
                 Inches(5.7), Inches(1.9), size=Pt(12), color=C_WHITE, wrap=True)

    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Market
# ═════════════════════════════════════════════════════════════════════════════

def slide_market(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Market", "Who needs this tool")

    segments = [
        ("OEM R&D Departments",
         ["Yamaha, Honda, KTM, Ducati chassis teams",
          "Prototype geometry validation before CAD",
          "Rapid what-if analysis for rake / trail / pivot height"],
         C_ACCENT),
        ("Racing Teams",
         ["MotoGP, WSBK, BSB — chassis setup engineers",
          "Anti-squat / anti-dive calibration per track",
          "Swingarm angle and sprocket ratio trade-offs"],
         C_CYAN),
        ("Suspension & Frame Manufacturers",
         ["Öhlins, WP, Bitubo — suspension tuning validation",
          "Custom frame builders: geometry optimisation",
          "Spring rate and sag calibration against load targets"],
         C_ACCENT2),
        ("Academia & Research",
         ["Motorcycle dynamics research groups",
          "Capstone / thesis projects on chassis design",
          "Teaching tool: Foale/Cossalter theory visualised live"],
         C_PURPLE),
        ("Advanced Enthusiasts & Tuners",
         ["High-end track day riders building custom bikes",
          "Aftermarket geometry kit validation (steering dampers, linkages)",
          "Self-build and café racer frame geometry checks"],
         C_WARN),
    ]

    col_w = Inches(2.4)
    for i, (title, items, color) in enumerate(segments):
        x = Inches(0.3 + i * 2.55)
        add_rect(s, x, Inches(1.25), col_w, Inches(5.9),
                 fill=RGBColor(0x16, 0x1c, 0x26), line=color, line_w=Pt(1.2))
        add_rect(s, x, Inches(1.25), col_w, Inches(0.05), fill=color)
        add_text(s, title, x + Inches(0.12), Inches(1.32),
                 col_w - Inches(0.2), Inches(0.55),
                 size=Pt(12), bold=True, color=color, wrap=True)
        txb = s.shapes.add_textbox(x + Inches(0.12), Inches(1.95),
                                   col_w - Inches(0.2), Inches(5.0))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"▸  {item}"
            run.font.size = Pt(10.5)
            run.font.color.rgb = C_WHITE

    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Reference Books
# ═════════════════════════════════════════════════════════════════════════════

def slide_references(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Reference Books & Formulas",
                "The physics behind every calculation")

    books = [
        {
            "title": "Motorcycle Handling\nand Chassis Design",
            "author": "Tony Foale  (2006)",
            "color": C_ACCENT,
            "chapters": [
                "Ch. 2 — Geometry: trail, rake, wheelbase",
                "Ch. 5 — Chassis types, swingarm angle",
                "Ch. 6 — Centre of gravity, weight distribution",
                "Ch. 8 — Anti-squat, anti-dive (graphical IC method)",
                "Ch. 9 — Suspension: spring rates, sag, damping",
                "Ch. 10 — Dynamics: load transfer, wheelie/stoppie",
                "Ch. 11 — Kinematics, chain wrap, WB change",
                "Ch. 13 — Aerodynamics: drag, lift, top speed",
            ],
        },
        {
            "title": "Motorcycle Dynamics\n2nd Edition",
            "author": "Vittore Cossalter  (2006)",
            "color": C_CYAN,
            "chapters": [
                "Ch. 1 — Chassis geometry, coordinate frames",
                "Ch. 2 — Tyre mechanics: contact patch, radii",
                "Ch. 4 — Aerodynamics: Cx, Cz, pitch moment",
                "Ch. 5 — Anti-squat R-ratio: tan(τ)/tan(σ)",
                "Ch. 6 — Fork compliance, trail change under braking",
                "Ch. 8 — Inertia properties, gyroscopic effects",
                "Ch. 5 — Squat ratio sweep vs sprocket teeth",
            ],
        },
        {
            "title": "Anti-Squat / Anti-Dive\nin Motorcycles (PDF)",
            "author": "Foale & Cossalter methods — §1–§9",
            "color": C_ACCENT2,
            "chapters": [
                "§1–§3 — Load transfer geometry, braking force",
                "§4 — Anti-dive brake caliper geometry",
                "§5 — Lean-angle correction to AS%",
                "§6 — Effective AD% with disc geometry",
                "§7 — AS% in cornering (lean sweep)",
                "§8 — Combined AS+AD under braking in corner",
                "§9 — Chain load contributions",
            ],
        },
    ]

    for i, b in enumerate(books):
        x = Inches(0.4 + i * 4.28)
        add_rect(s, x, Inches(1.3), Inches(4.0), Inches(5.8),
                 fill=RGBColor(0x14, 0x1a, 0x23), line=b["color"], line_w=Pt(1.5))
        add_rect(s, x, Inches(1.3), Inches(4.0), Inches(0.07), fill=b["color"])

        add_text(s, b["title"], x + Inches(0.18), Inches(1.42),
                 Inches(3.7), Inches(0.72),
                 size=Pt(14), bold=True, color=b["color"], wrap=True)
        add_text(s, b["author"], x + Inches(0.18), Inches(2.15),
                 Inches(3.7), Inches(0.35),
                 size=Pt(11), color=C_MUTED)

        txb = s.shapes.add_textbox(x + Inches(0.18), Inches(2.52),
                                   Inches(3.7), Inches(4.3))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for ch in b["chapters"]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(5)
            run = p.add_run()
            run.text = f"→  {ch}"
            run.font.size = Pt(10.5)
            run.font.color.rgb = C_WHITE

    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Coverage Map
# ═════════════════════════════════════════════════════════════════════════════

def slide_coverage(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Physics Coverage",
                "13 modules — what the software computes")

    modules = [
        # (name, detail, color)
        ("Chassis Geometry",     "Trail · rake · swingarm angle · WB",          C_ACCENT),
        ("Centre of Gravity",    "X_cg · Y_cg · F/R weight split",              C_ACCENT),
        ("Anti-Squat / IC",      "Foale AS% + Cossalter R-ratio",               C_WARN),
        ("Anti-Dive",            "Brake caliper geometry · lean correction",     C_WARN),
        ("Suspension Dynamics",  "WR · nat freq · sag% · ζ · optimal damping",  C_CYAN),
        ("Ergonomics",           "Knee/hip angles · forward lean · reach",       C_ACCENT2),
        ("Load Transfer",        "Braking/accel/cornering · wheelie/stoppie",    C_ACCENT),
        ("Tire Physics",         "Contact patch · deflection · dynamic radius",  C_CYAN),
        ("Kinematics",           "Rear axle locus · WB change · chain wrap",     C_ACCENT2),
        ("Inertia",              "I_yaw · I_pitch · radii of gyration",          C_PURPLE),
        ("Fork Compliance",      "Deflection · trail change · SAT under braking",C_WARN),
        ("Aerodynamics",         "Drag · lift · pitch moment · gear-ltd V_max",  C_ACCENT),
        ("Stability / Handling", "SI · AI · wobble sensitivity · lean limit",    C_ACCENT2),
    ]

    cols = 4
    rows = 4
    cw = Inches(3.1)
    ch = Inches(1.35)
    pad_x = Inches(0.22)
    pad_y = Inches(0.25)
    start_x = Inches(0.3)
    start_y = Inches(1.3)

    for idx, (name, detail, color) in enumerate(modules):
        r = idx // cols
        c = idx % cols
        x = start_x + c * (cw + pad_x)
        y = start_y + r * (ch + pad_y)
        add_rect(s, x, y, cw, ch,
                 fill=RGBColor(0x16, 0x1c, 0x26), line=color, line_w=Pt(1.0))
        add_rect(s, x, y, Inches(0.06), ch, fill=color)
        add_text(s, name, x + Inches(0.18), y + Inches(0.12),
                 cw - Inches(0.25), Inches(0.42),
                 size=Pt(12), bold=True, color=color)
        add_text(s, detail, x + Inches(0.18), y + Inches(0.55),
                 cw - Inches(0.25), Inches(0.72),
                 size=Pt(10), color=C_MUTED, wrap=True)

    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Technical Stack
# ═════════════════════════════════════════════════════════════════════════════

def slide_stack(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_title(s, "Technical Architecture",
                "How the software is built")

    layers = [
        ("FRONTEND",
         "Electron + React 18 + Vite + TypeScript",
         "22 analysis tabs  ·  Zustand state  ·  Recharts  ·  Three.js 3D  ·  "
         "Instant TypeScript physics (150ms debounce before backend sync)  ·  "
         "ChassisViz2D live engineering drawing",
         C_ACCENT, Inches(1.3)),
        ("BACKEND",
         "Python FastAPI  ·  uvicorn  ·  auto-started by Electron",
         "POST /api/dag-analysis → 13 physics modules in one call  ·  "
         "DAG topological solver (40+ coupled parameters)  ·  "
         "6 standalone engine modules: suspension, ergonomics, tire, kinematics, fork, aero  ·  "
         "Validated at 86% against published specs",
         C_CYAN, Inches(2.85)),
        ("PHYSICS ENGINE",
         "Foale Eq 5.1–11.x  ·  Cossalter Eq 1–8  ·  520-chain geometry",
         "Tangent contact point (not CS center) for IC computation  ·  "
         "atan2 swingarm angle  ·  Gear-limited top speed: min(P-limited, RPM-ceiling)  ·  "
         "Cossalter R = tan(τ)/tan(σ) from rear contact patch  ·  "
         "8 bike family presets with calibrated spring rates",
         C_ACCENT2, Inches(4.4)),
        ("SIMULATION",
         "Multi-Body Dynamics (MBD)  ·  Phases 1–3 complete (32/32 tests)",
         "Generalized-α integrator  ·  Rigid body constraints  ·  "
         "Phase 4 next: GJK collision, LCP contact, Pacejka tire model  ·  "
         "Validated chassis_sim quasi-static engine",
         C_PURPLE, Inches(5.95)),
    ]

    for label, tech, detail, color, y in layers:
        add_rect(s, Inches(0.4), y, Inches(12.5), Inches(1.2),
                 fill=RGBColor(0x16, 0x1c, 0x26), line=color, line_w=Pt(1.0))
        add_rect(s, Inches(0.4), y, Inches(0.07), Inches(1.2), fill=color)
        add_text(s, label, Inches(0.65), y + Inches(0.08),
                 Inches(1.5), Inches(0.4), size=Pt(10), bold=True, color=color)
        add_text(s, tech, Inches(2.2), y + Inches(0.06),
                 Inches(4.2), Inches(0.42), size=Pt(11), bold=True, color=C_WHITE)
        add_text(s, detail, Inches(2.2), y + Inches(0.5),
                 Inches(10.5), Inches(0.65), size=Pt(10), color=C_MUTED, wrap=True)

    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Closing
# ═════════════════════════════════════════════════════════════════════════════

def slide_close(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, Inches(0), Inches(0), W, H, fill=RGBColor(0x0a, 0x0e, 0x14))
    add_rect(s, Inches(0), Inches(0), Inches(0.18), H, fill=C_ACCENT)
    add_rect(s, Inches(0), H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)

    add_text(s, "Where We Are Today",
             Inches(0.55), Inches(0.9), Inches(10), Inches(0.7),
             size=Pt(30), bold=True, color=C_WHITE)

    stats = [
        ("22", "Analysis Tabs", C_ACCENT),
        ("13", "Physics Modules", C_CYAN),
        ("8",  "Bike Presets", C_ACCENT2),
        ("86%","Validation Pass", C_WARN),
        ("38","Formula Checks",  C_PURPLE),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = Inches(0.55 + i * 2.55)
        add_rect(s, x, Inches(1.85), Inches(2.3), Inches(1.5),
                 fill=RGBColor(0x16, 0x1c, 0x26), line=color, line_w=Pt(1.5))
        add_text(s, num, x + Inches(0.15), Inches(1.92),
                 Inches(2.0), Inches(0.7), size=Pt(36), bold=True, color=color)
        add_text(s, label, x + Inches(0.15), Inches(2.65),
                 Inches(2.0), Inches(0.55), size=Pt(11), color=C_MUTED)

    add_text(s, "What's Next",
             Inches(0.55), Inches(3.65), Inches(10), Inches(0.5),
             size=Pt(20), bold=True, color=C_ACCENT)

    nexts = [
        "Fix /api/sweep endpoint — motion ratio MR (currently blows up on default shock geometry)",
        "MBD Phase 4 — GJK collision detection + LCP contact solver + tire contact",
        "Pacejka magic formula tire model (Phase 12) for combined slip scenarios",
        "Whipple eigenvalue modes vs speed — capsize / weave / wobble stability analysis",
        "Damping coefficient UI sliders + kinematics positions locus visualization",
    ]
    txb = s.shapes.add_textbox(Inches(0.55), Inches(4.2), Inches(12.3), Inches(2.8))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in nexts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(12)
        run.font.color.rgb = C_WHITE


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    prs = new_prs()

    slide_hero(prs)
    slide_background(prs)
    slide_aim(prs)
    slide_market(prs)
    slide_references(prs)
    slide_coverage(prs)
    slide_stack(prs)
    slide_close(prs)

    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "chassis_workbench_presentation.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
